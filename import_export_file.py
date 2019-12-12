import pandas as pd
import os
import pptx

def get_path(file_name,in_or_out):
    # ______________________________________________________________________________________________
    # DESCRIPTION

    # This function generates the file and folder path locations
    # to call it use get_path(file_name, in_or_out = input or output)
    # ______________________________________________________________________________________________

    # ______________________________________________________________________________________________
    # GENERATING STRINGS
    # Get current path
    main_path = os.getcwd()

    # Generate file & path strings
    if in_or_out == 'input':
        folder_path = main_path + '\inputs'
        file_path = main_path + '\inputs' + '\\' + file_name
    if in_or_out == 'output':
        folder_path = main_path + '\outputs'
        file_path = main_path + '\outputs' + '\\' + file_name

    return folder_path, file_path
    # ______________________________________________________________________________________________

def export_data(table, file_name,state='on'):
    # __________________________________________________________________________________________________________________
    # DESCRIPTION

    # This function exports files to the folder "outputs"
    # to call it use export_data(table, file_name)
    # ------------------------------------------------------------------------------------------------------------------
    # ARGUMENTS

    # table
    # file_name
    # state
    # __________________________________________________________________________________________________________________

    # Enable/disable feature
    if state=='off':
        return
    # ------------------------------------------------------------------------------------------------------------------
    # EXPORTING EXCEL OUTPUT

    # Generate file & folder path strings
    output_path, output_file = get_path(file_name, 'output')

    # Export file
    table.to_excel(output_file)

    # Print success
    #print('The file '+file_name+' was successfully exported to '+output_path+'\n')

    return

def import_data(file_name):
    # ______________________________________________________________________________________________
    # DESCRIPTION

    # This function imports files from the folder inputs
    # to call it use import_data(file_name)
    # ______________________________________________________________________________________________

    # ______________________________________________________________________________________________
    # IMPORTING FILE

    # Generate file & folder path strings
    input_path, input_file = get_path(file_name,'input')

    # Handle error
    if file_name in os.listdir(input_path):
        print('The input file ' +file_name +' was imported successfully. \n')
    else:
        print('Could not find input file ' +file_name +' in the folder '+input_path)
        return 0

    # Import file
    df = pd.read_excel(input_file)

    return df
    # ______________________________________________________________________________________________

def export_nice_excel(volume_sheet,finance_sheet,complete_sales_history,no_of_months):
    # MAKE TABLE LOOK NICE

    # Set paths
    main_path = os.getcwd()
    file_path_format = main_path + '\outputs' + '\\' + 'AtomBrands_financials.xlsx'

    # Get number of rows
    number_rows = len(volume_sheet.index)

    # Create custom writer
    writer = pd.ExcelWriter(file_path_format,
                            engine='xlsxwriter')
    finance_sheet.to_excel(writer,
                           index=False,
                           float_format="%.2f",
                           sheet_name='Price Tree',
                           startrow=1)
    volume_sheet.to_excel(writer,
                          index=False,
                          float_format="%.2f",
                          sheet_name='Volume Breakdown',
                          startrow=1)
    complete_sales_history.to_excel(writer,
                                    index=False,
                                    float_format="%.2f",
                                    sheet_name='Complete Sales History',
                                    startrow=0)

    # Set excel writer config
    workbook = writer.book
    worksheet = writer.sheets['Price Tree']
    worksheet2 = writer.sheets['Volume Breakdown']

    # Add a money formats for currency cells
    # RMB
    rmb_fmt = workbook.add_format({'num_format': '¥#,##0.00', 'bold': False, 'align': 'right'})
    rmb_fmt_border = workbook.add_format({'num_format': '¥#,##0.00', 'bold': False, 'align': 'right'})
    rmb_fmt_border.set_right(1)
    # GBP
    gbp_fmt = workbook.add_format({'num_format': '£#,##0.00', 'bold': False, 'align': 'right'})
    gbp_fmt_border = workbook.add_format({'num_format': '£#,##0.00', 'bold': False, 'align': 'right'})
    gbp_fmt_border.set_right(1)

    # Add a percent format with 2 decimal points for % cells
    percent_fmt = workbook.add_format({'num_format': '0.00%', 'bold': False, 'align': 'center'})
    percent_fmt.set_right(1)

    # Add other formats
    other_fmt = workbook.add_format({'bold': False, 'align': 'center'})

    # Add title color formats
    col_title1 = workbook.add_format({'bold': True, 'font_color': '#ffffff', 'bg_color': '#008484', 'align': 'center'})
    col_title2 = workbook.add_format({'bold': True, 'font_color': '#008484', 'bg_color': '#ffffff', 'align': 'center'})
    col_title1.set_right(1)
    col_title1.set_right(1)

    # Add border formats
    right_border = workbook.add_format({'bold': False, 'align': 'center'})
    right_border.set_right(1)

    # PRICE TREE SHEET CONFIG
    # Set Titles
    worksheet.merge_range('A1:C1', 'PRODUCT ID', col_title1)
    worksheet.merge_range('D1:H1', 'PRODUCT INFO', col_title2)
    worksheet.merge_range('I1:O1', 'EXW, VILC & MARGIN', col_title1)
    worksheet.merge_range('P1:Y1', 'RRP & NR', col_title2)
    worksheet.merge_range('Z1:AH1', 'COSTS', col_title1)
    worksheet.merge_range('AI1:AP1', 'FLOATING MACOS', col_title2)

    # Set column width and data format
    worksheet.set_column('A:B', 15, other_fmt)
    worksheet.set_column('C:C', 18, right_border)
    worksheet.set_column('D:D', 16)
    worksheet.set_column('E:E', 21)
    worksheet.set_column('F:F', 65)
    worksheet.set_column('G:G', 5, other_fmt)
    worksheet.set_column('H:H', 5, right_border)
    worksheet.set_column('I:I', 18, gbp_fmt)
    worksheet.set_column('J:J', 18, rmb_fmt)
    worksheet.set_column('K:K', 18, gbp_fmt)
    worksheet.set_column('L:L', 18, rmb_fmt)
    worksheet.set_column('M:M', 18, gbp_fmt)
    worksheet.set_column('N:N', 18, rmb_fmt)
    worksheet.set_column('O:O', 10, percent_fmt)
    worksheet.set_column('P:P', 18, gbp_fmt)
    worksheet.set_column('Q:Q', 18, rmb_fmt)
    worksheet.set_column('R:R', 18, gbp_fmt)
    worksheet.set_column('S:S', 18, rmb_fmt)
    worksheet.set_column('T:T', 18, gbp_fmt)
    worksheet.set_column('U:U', 18, rmb_fmt)
    worksheet.set_column('V:V', 18, gbp_fmt)
    worksheet.set_column('W:W', 18, rmb_fmt)
    worksheet.set_column('X:X', 18, gbp_fmt)
    worksheet.set_column('Y:Y', 18, rmb_fmt_border)
    worksheet.set_column('Z:AF', 18, rmb_fmt)
    worksheet.set_column('AG:AG', 18, gbp_fmt)
    worksheet.set_column('AH:AH', 18, rmb_fmt_border)
    worksheet.set_column('AI:AI', 18, gbp_fmt)
    worksheet.set_column('AJ:AJ', 18, rmb_fmt)
    worksheet.set_column('AK:AK', 18, gbp_fmt)
    worksheet.set_column('AL:AL', 18, rmb_fmt)
    worksheet.set_column('AM:AM', 18, gbp_fmt)
    worksheet.set_column('AN:AN', 18, rmb_fmt)
    worksheet.set_column('AO:AO', 18, gbp_fmt)
    worksheet.set_column('AP:AP', 18, rmb_fmt_border)

    # VOLUME BREAKDOWN CONFIG
    # Set Titles
    end_month_title_col = 26 + no_of_months

    worksheet2.merge_range('A1:D1', 'PRODUCT INFO', col_title1)
    worksheet2.merge_range('E1:V1', 'REVENUE & MACO', col_title2)
    worksheet2.merge_range('W1:AA1', 'VOLUME PER CHANNEL', col_title1)
    worksheet2.merge_range(0, 27, 0, end_month_title_col, 'VOLUME PER MONTH', col_title2)
    worksheet2.merge_range(0, end_month_title_col + 1, 0, end_month_title_col + 10, 'TOTAL REVENUE BY CHANNEL',
                           col_title1)
    worksheet2.merge_range(0, end_month_title_col + 11, 0, end_month_title_col + 20, 'TOTAL PROFIT BY CHANNEL',
                           col_title2)
    worksheet2.merge_range(0, end_month_title_col + 21, 0, end_month_title_col + 20 + (no_of_months * 2),
                           'TOTAL REVENUE BY MONTH', col_title1)
    worksheet2.merge_range(0, end_month_title_col + 20 + (no_of_months * 2) + 1, 0,
                           end_month_title_col + 20 + (no_of_months * 4)
                           , 'TOTAL PROFIT BY MONTH', col_title2)

    # Set column width and data format
    worksheet2.set_column('A:B', 15, other_fmt)
    worksheet2.set_column('C:C', 21)
    worksheet2.set_column('D:D', 5, right_border)
    worksheet2.set_column('E:E', 18, gbp_fmt)
    worksheet2.set_column('F:F', 18, rmb_fmt)
    worksheet2.set_column('G:G', 18, gbp_fmt)
    worksheet2.set_column('H:H', 18, rmb_fmt)
    worksheet2.set_column('I:I', 18, gbp_fmt)
    worksheet2.set_column('J:J', 18, rmb_fmt)
    worksheet2.set_column('K:K', 18, gbp_fmt)
    worksheet2.set_column('L:L', 18, rmb_fmt)
    worksheet2.set_column('M:M', 18, gbp_fmt)
    worksheet2.set_column('N:N', 18, rmb_fmt)
    worksheet2.set_column('O:O', 18, gbp_fmt)
    worksheet2.set_column('P:P', 18, rmb_fmt)
    worksheet2.set_column('Q:Q', 18, gbp_fmt)
    worksheet2.set_column('R:R', 18, rmb_fmt)
    worksheet2.set_column('S:S', 18, gbp_fmt)
    worksheet2.set_column('T:T', 18, rmb_fmt)
    worksheet2.set_column('U:U', 18, gbp_fmt)
    worksheet2.set_column('V:V', 18, rmb_fmt_border)
    worksheet2.set_column('W:Z', 10, other_fmt)
    worksheet2.set_column('AA:AA', 10, right_border)
    worksheet2.set_column(27, end_month_title_col - 1, 10, other_fmt)
    worksheet2.set_column(end_month_title_col, end_month_title_col, 10, right_border)
    worksheet2.set_column(end_month_title_col + 1, end_month_title_col + 1, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 2, end_month_title_col + 2, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 3, end_month_title_col + 3, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 4, end_month_title_col + 4, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 5, end_month_title_col + 5, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 6, end_month_title_col + 6, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 7, end_month_title_col + 7, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 8, end_month_title_col + 8, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 9, end_month_title_col + 9, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 10, end_month_title_col + 10, 18, rmb_fmt_border)
    worksheet2.set_column(end_month_title_col + 11, end_month_title_col + 11, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 12, end_month_title_col + 12, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 13, end_month_title_col + 13, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 14, end_month_title_col + 14, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 15, end_month_title_col + 15, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 16, end_month_title_col + 16, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 17, end_month_title_col + 17, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 18, end_month_title_col + 18, 18, rmb_fmt)
    worksheet2.set_column(end_month_title_col + 19, end_month_title_col + 19, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 20, end_month_title_col + 20, 18, rmb_fmt_border)

    i = 1
    j = 1
    while i < no_of_months:
        worksheet2.set_column(end_month_title_col + 20 + j, end_month_title_col + 20 + j, 18, gbp_fmt)
        worksheet2.set_column(end_month_title_col + 21 + j, end_month_title_col + 21 + j, 18, rmb_fmt)
        i = i + 1
        j = j + 2
    worksheet2.set_column(end_month_title_col + 20 + j, end_month_title_col + 20 + j, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 21 + j, end_month_title_col + 21 + j, 18, rmb_fmt_border)

    j = j + 2
    i = 1
    while i < no_of_months:
        worksheet2.set_column(end_month_title_col + 20 + j, end_month_title_col + 20 + j, 18, gbp_fmt)
        worksheet2.set_column(end_month_title_col + 21 + j, end_month_title_col + 21 + j, 18, rmb_fmt)
        i = i + 1
        j = j + 2
    worksheet2.set_column(end_month_title_col + 20 + j, end_month_title_col + 20 + j, 18, gbp_fmt)
    worksheet2.set_column(end_month_title_col + 21 + j, end_month_title_col + 21 + j, 18, rmb_fmt_border)

    writer.save()

    return file_path_format

    return

def export_ppt():

    main_path = os.getcwd()
    file_path_format = main_path + '\outputs' + '\\' + 'test.pptx'

    prs = pptx.Presentation()
    prs.save(file_path_format)

    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "TEST TITLE"

    return